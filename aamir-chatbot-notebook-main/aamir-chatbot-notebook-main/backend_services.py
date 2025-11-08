import os
import base64
import requests
import tempfile
from bs4 import BeautifulSoup
import PyPDF2
import docx2txt
from pptx import Presentation
import networkx as nx
from gtts import gTTS
from io import BytesIO

from pptx import Presentation
from langchain_core.documents import Document
from dotenv import load_dotenv
import os
from langchain_google_genai import ChatGoogleGenerativeAI , GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.tools.tavily_search import TavilySearchResults
from langchain_core.messages import HumanMessage
# from langchain.embeddings import OpenAIEmbeddings

# Global variable to store the vector database
vector_db = None

load_dotenv()

embeddings = GoogleGenerativeAIEmbeddings(model="models/text-embedding-004")
llm = ChatGoogleGenerativeAI(model="models/gemini-2.5-flash")



def extract_ppt(file_path):
    """Loads a PPTX file using python-pptx and converts it to a list of LangChain Documents."""
    prs = Presentation(file_path)
    langchain_documents = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_title = f"Slide {i + 1}"  # Default title
        
        # 1. Safely check for a title placeholder (Placeholder idx 0 is usually the title)
        title_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                title_shape = shape
                break

        if title_shape and hasattr(title_shape, 'text') and title_shape.text:
            slide_title = title_shape.text
            
        # 2. Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text
                if text:
                    slide_text.append(text)
        
        # 3. Create a Document for each slide
        doc = Document(
            page_content="\n".join(slide_text),
            metadata={
                "source": file_path,
                "slide_number": i + 1, 
                "title": slide_title
            }
        )
        langchain_documents.append(doc)

    return langchain_documents

from langchain_community.document_loaders import Docx2txtLoader

def extract_word(file_path):
    loader = Docx2txtLoader(file_path)
    word_documents = loader.load()
    
    return word_documents

from langchain_community.document_loaders import PyPDFLoader

def extract_pdf(file_path):
    loader=PyPDFLoader(file_path)
    docs = loader.load()
    print(docs)


def image_to_base64(image_path):
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode('utf-8')

def extract_text_from_file(file):
    """Extract text from various file types including PDF, DOCX, PPTX, and images."""
    documents = []
    file_extension = file.name.split('.')[-1].lower()
    
    try:
        # Create a temporary file to save the uploaded file content
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as tmp_file:
            # Write the uploaded file's content to the temporary file
            tmp_file.write(file.getvalue())
            tmp_file_path = tmp_file.name
        
        try:
            if file_extension == 'pdf':
                from langchain_community.document_loaders import PyPDFLoader
                loader = PyPDFLoader(tmp_file_path)
                documents = loader.load()
                # Update metadata to use original filename
                for doc in documents:
                    doc.metadata['source'] = file.name
                
            elif file_extension == 'docx':
                from langchain_community.document_loaders import Docx2txtLoader
                loader = Docx2txtLoader(tmp_file_path)
                documents = loader.load()
                # Update metadata to use original filename
                for doc in documents:
                    doc.metadata['source'] = file.name
                
            elif file_extension in ['pptx', 'ppt']:
                prs = Presentation(tmp_file_path)
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        doc = Document(
                            page_content="\n".join(slide_text),
                            metadata={
                                "source": file.name,  # Already using original filename
                                "slide_number": i + 1
                            }
                        )
                        documents.append(doc)
                        
            elif file_extension in ['jpg', 'jpeg', 'png']:
                # For images, we'll just store the filename for now
                # In a real RAG system, you'd use OCR or a vision model here
                doc = Document(
                    page_content=f"[Image: {file.name}]",
                    metadata={"source": file.name, "type": "image"}  # Already using original filename
                )
                documents.append(doc)
                
        finally:
            # Clean up the temporary file
            try:
                os.unlink(tmp_file_path)
            except:
                pass
                
    except Exception as e:
        raise Exception(f"Error processing {file.name}: {str(e)}")
        
    return documents

from langchain_community.document_loaders import WebBaseLoader
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnableLambda
import bs4

def extract_text_from_url(url):
    """
    Extract text content from a given URL using WebBaseLoader.
    
    Args:
        url (str): The URL to extract text from
        
    Returns:
        list: List of Document objects containing the extracted text
        
    Raises:
        Exception: If there's an error loading or processing the URL
    """
    try:
        # Validate URL format
        if not url.startswith(('http://', 'https://')):
            url = f'https://{url}'
            
        # Configure WebBaseLoader with custom settings
        loader = WebBaseLoader(
            web_path=url,
            # Add any additional configuration here if needed
            # e.g., continue_on_failure=True,
            #       requests_per_second=2,
        )
        
        # Load and return documents
        return loader.load()
        
    except Exception as e:
        raise Exception(f"Error processing URL {url}: {str(e)}")

from dotenv import load_dotenv


def process_documents(uploaded_files, website_urls=None):
    """
    Process all uploaded files and website URLs to create a vector database.
    
    Args:
        uploaded_files: List of uploaded file objects
        website_urls: List of website URLs (optional)
        
    Returns:
        tuple: (success: bool, message: str, db: Chroma)
    """
    try:
        all_documents = []
        
        # Process uploaded files
        for file in uploaded_files:
            try:
                docs = extract_text_from_file(file)
                if docs:
                    all_documents.extend(docs)
            except Exception as e:
                print(f"Error processing file {file.name}: {str(e)}")
                continue
        
        # Process website URLs
        if website_urls:
            for url in website_urls:
                try:
                    docs = extract_text_from_url(url)
                    if docs:
                        all_documents.extend(docs)
                except Exception as e:
                    print(f"Error processing URL {url}: {str(e)}")
                    continue
        
        if not all_documents:
            return False, "No valid documents or text content found to process.", None
        
        # Split documents into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            is_separator_regex=False,
        )
        chunks = text_splitter.split_documents(all_documents)
        
        # Create in-memory vector store
        global vector_db
        vector_db = Chroma.from_documents(
            documents=chunks,
            embedding=embeddings,
            persist_directory=None  # This ensures it's not persisted to disk
        )
        
        return True, f"Successfully processed {len(chunks)} document chunks.", vector_db
        
    except Exception as e:
        return False, f"Error processing documents: {str(e)}", None



def generate_knowledge_graph(texts):
    """Generate a knowledge graph from the provided texts."""
    # Simple keyword extraction and relationship mapping
    # In a real implementation, you'd use NLP techniques
    G = nx.Graph()
    
    # Add nodes and edges based on co-occurrence
    for i, text in enumerate(texts):
        # Simple tokenization (in reality, use NLP for better results)
        words = [word.lower() for word in text.split() if len(word) > 3][:20]  # Limit words for demo
        
        # Add nodes
        for word in words:
            if word not in G:
                G.add_node(word, title=word, group=1)
        
        # Add edges between co-occurring words
        for j in range(len(words)):
            for k in range(j + 1, len(words)):
                if G.has_edge(words[j], words[k]):
                    G[words[j]][words[k]]['weight'] += 1
                else:
                    G.add_edge(words[j], words[k], weight=1)
    
    return G

def text_to_speech(text, max_retries=2):
    """
    Convert text to speech and return the audio data as base64.
    
    Args:
        text (str): The text to convert to speech
        max_retries (int): Number of retry attempts if the request fails
        
    Returns:
        str: Base64 encoded audio data or None if conversion fails
    """
    if not text or not text.strip():
        print("Error: Empty or invalid text provided for speech synthesis")
        return None
    
    # Truncate very long text to avoid API issues
    if len(text) > 5000:
        text = text[:5000] + "... [text truncated]"
    
    # Try multiple times in case of temporary network issues
    for attempt in range(max_retries + 1):
        try:
            print(f"Attempt {attempt + 1}: Generating speech for text (first 100 chars): {text[:100]}...")
            
            # Try with a timeout to prevent hanging
            tts = gTTS(text=text, lang='en', slow=False)
            
            # Use a buffer to store the audio data
            audio_buffer = BytesIO()
            tts.write_to_fp(audio_buffer)
            audio_buffer.seek(0)
            audio_data = audio_buffer.read()
            
            if not audio_data:
                print("Error: No audio data was generated")
                continue  # Try again if no data
                
            audio_base64 = base64.b64encode(audio_data).decode('utf-8')
            print(f"Successfully generated audio (base64 length: {len(audio_base64)})")
            return audio_base64
            
        except Exception as e:
            error_msg = f"Attempt {attempt + 1} failed: {str(e)}"
            print(error_msg)
            
            # If we've exhausted all retries, try a fallback
            if attempt == max_retries:
                print("All attempts failed. Falling back to simple notification.")
                # Return None to indicate failure
                return None
                
            # Wait a bit before retrying
            import time
            time.sleep(1)
    
    return None

def summarize_documents(docs):
    """
    Summarize documents using a map-reduce approach.
    
    Args:
        docs: List of Document objects to be summarized
        
    Returns:
        str: A comprehensive summary of all documents
    """
    # 1. Split documents into chunks
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=500)
    chunks = text_splitter.split_documents(docs)
    
    # 2. Define the MAP Step
    map_prompt = PromptTemplate.from_template(
        "Write a concise summary of the following chunk of text:\n{context}\nCONCISE SUMMARY:"
    )
    map_chain = map_prompt | llm | StrOutputParser()
    
    # 3. Define the REDUCE Step
    reduce_prompt = PromptTemplate.from_template(
        "The following are concise summaries. Combine them into a single, comprehensive final summary:\n{context}\nFINAL SUMMARY:"
    )
    reduce_chain = reduce_prompt | llm | StrOutputParser()
    
    # 4. Map-Reduce Process
    # MAP: Apply the map_chain to every document chunk
    summaries = [map_chain.invoke({"context": chunk.page_content}) for chunk in chunks]
    
    # REDUCE: Combine all the summaries into a single string for the final prompt
    combined_summaries = "\n\n---\n\n".join(summaries)
    
    # FINAL REDUCE CALL
    final_summary = reduce_chain.invoke({"context": combined_summaries})
    
    return final_summary


def get_mock_response(user_input, db=None):
    """Generate a response using RAG with the vector database.
    
    Args:
        user_input (str): The user's query
        db: Optional vector database instance. If None, uses the global vector_db
        
    Returns:
        dict: Dictionary containing response text and sources
            {
                'response': str,  # The generated response
                'sources': list   # List of source documents with metadata
            }
    """
    try:
        # Get relevant documents from the database
        if db is None:
            from app import vector_db
            db = vector_db
            
        if db is None:
            return {
                'response': "I don't have any knowledge base to reference. Please upload some documents first.",
                'sources': []
            }
            
        # Get the most relevant documents
        docs = db.similarity_search(user_input, k=3)
        
        if not docs:
            return {
                'response': "I couldn't find any relevant information in the knowledge base to answer your question.",
                'sources': []
            }
            
        # Extract unique sources from documents
        sources = []
        seen_sources = set()
        
        for doc in docs:
            source = doc.metadata.get('source', 'Unknown Source')
            if source not in seen_sources:
                seen_sources.add(source)
                sources.append({
                    'source': source,
                    'page_content': doc.page_content[:200] + '...' if len(doc.page_content) > 200 else doc.page_content
                })
        
        # Format the context from the documents
        context = "\n\n".join([f"Document {i+1}:\n{doc.page_content}" for i, doc in enumerate(docs)])
        
        # Create a prompt template
        template = """Use the following pieces of context to answer the question at the end. 
        If you don't know the answer, just say that you don't know, don't try to make up an answer.
        
        {context}
        
        Question: {question}
        
        helpful Answer:"""
        
        # Generate response using the LLM
        from langchain_core.prompts import PromptTemplate
        # from langchain_openai import ChatOpenAI
        
        prompt = PromptTemplate(template=template, input_variables=["context", "question"])
        embeddings = GoogleGenerativeAIEmbeddings(model="models/text-embedding-004")
        llm = ChatGoogleGenerativeAI(model="models/gemini-2.5-flash")
        
        chain = prompt | llm
        
        # Generate the response
        response = chain.invoke({"context": context, "question": user_input})
        
        return {
            'response': response.content,
            'sources': sources
        }
        
    except Exception as e:
        error_msg = f"Error generating response: {str(e)}"
        print(error_msg)  # Log the error for debugging
        return "I encountered an error while processing your request. Please try again later."


def get_tavily_search_response(user_query):
    """
    Generate a response using Tavily search tool and LLM with LCEL.
    
    Args:
        user_query (str): The user's search query
        
    Returns:
        str: Generated response based on Tavily search results
    """
    try:
        # Initialize Tavily search tool
        tavily_tool = TavilySearchResults(
            max_results=5,
            search_depth="advanced",
            include_answer=True,
            include_raw_content=False,
            include_images=False
        )
        
        # Bind the tool to the LLM
        llm_with_tools = llm.bind_tools([tavily_tool])
        
        # Get search results from Tavily
        print(f"Searching for: {user_query}")
        search_results = tavily_tool.invoke({"query": user_query})
        
        # Format search results into a readable context
        if not search_results:
            return "I couldn't find any relevant information from the web. Please try rephrasing your query."
        
        # Extract and format the search results
        context_parts = []
        for idx, result in enumerate(search_results, 1):
            if isinstance(result, dict):
                title = result.get('title', 'No title')
                content = result.get('content', result.get('snippet', 'No content'))
                url = result.get('url', '')
                context_parts.append(f"{idx}. {title}\n{content}\nSource: {url}\n")
        
        search_context = "\n".join(context_parts)
        
        # Create a prompt template for generating response using LCEL
        prompt_template = """You are a helpful AI assistant with access to web search results. 
Use the following search results to provide a comprehensive and accurate answer to the user's question.
Cite the sources when relevant.

Search Results:
{search_context}

User Question: {user_query}

Provide a detailed and well-structured answer based on the search results above:"""
        
        prompt = PromptTemplate(
            template=prompt_template,
            input_variables=["search_context", "user_query"]
        )
        
        # Create the LCEL chain
        chain = (
            {"search_context": lambda x: search_context, "user_query": lambda x: x}
            | prompt
            | llm
            | StrOutputParser()
        )
        
        # Generate the final response
        response = chain.invoke(user_query)
        
        return response
        
    except Exception as e:
        error_msg = f"Error with Tavily search: {str(e)}"
        print(error_msg)
        return f"I encountered an error while searching the web: {str(e)}. Please try again later."
    